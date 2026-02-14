"""
CAPABILITIES SHOWCASE PRESENTATION
==================================
Demonstrates all programmatic PowerPoint generation capabilities.

This script generates an 8-slide presentation showcasing:
1. Title slide with custom design
2. Chart generation via QuickChart.io
3. Chart animation with morph transitions
4. AI-generated images via Ideogram
5. Stock photography integration
6. Diagram generation via Mermaid
7. Avatar grid with DiceBear
8. Full integration (all capabilities combined)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

# Import our API utilities
from api_utils import APIUtils


class PresentationGenerator:
    """Main presentation generator class"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)

        # Create assets directory
        self.assets_dir = "assets"
        os.makedirs(self.assets_dir, exist_ok=True)

        # Color scheme
        self.colors = {
            'primary_blue': RGBColor(0, 120, 212),      # #0078d4
            'accent_purple': RGBColor(135, 100, 184),   # #8764b8
            'dark_gray': RGBColor(50, 49, 48),          # #323130
            'medium_gray': RGBColor(96, 94, 92),        # #605e5c
            'light_gray': RGBColor(243, 242, 241),      # #f3f2f1
            'white': RGBColor(255, 255, 255),
            'success': RGBColor(16, 124, 16),           # #107c10
        }

    def add_morph_transition(self, slide, duration_ms=1000):
        """Add morph transition to a slide"""
        xml = f'''
        <mc:AlternateContent
          xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
          <mc:Choice
            xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main"
            Requires="p159">
            <p:transition
              xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
              spd="slow" p14:dur="{duration_ms}">
              <p159:morph option="byObject"/>
            </p:transition>
          </mc:Choice>
          <mc:Fallback>
            <p:transition
              xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
              spd="slow">
              <p:fade/>
            </p:transition>
          </mc:Fallback>
        </mc:AlternateContent>
        '''
        slide.element.append(etree.fromstring(xml))

    def add_fade_transition(self, slide, duration_ms=800):
        """Add fade transition to a slide"""
        xml = f'''
        <mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
          <mc:Choice xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" Requires="p14">
            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          spd="slow" p14:dur="{duration_ms}">
              <p:fade />
            </p:transition>
          </mc:Choice>
          <mc:Fallback>
            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
              <p:fade />
            </p:transition>
          </mc:Fallback>
        </mc:AlternateContent>
        '''
        slide.element.append(etree.fromstring(xml))

    def slide_1_title(self):
        """Slide 1: Title slide with custom design"""
        print("Creating Slide 1: Title Slide...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Background gradient (simulated with shape)
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors['light_gray']
        background.line.color.rgb = self.colors['light_gray']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Programmatic PowerPoint\nCapabilities Showcase"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(54)
                run.font.bold = True
                run.font.color.rgb = self.colors['primary_blue']
                run.font.name = 'Segoe UI'

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(11.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Demonstrating python-pptx + API Integrations"
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for run in subtitle_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.color.rgb = self.colors['medium_gray']
            run.font.name = 'Segoe UI'

        # Generate and add DiceBear avatar
        avatar_path = f"{self.assets_dir}/avatar_speaker.png"
        if APIUtils.get_dicebear_avatar("showcase", avatar_path, style="bottts", size=150):
            slide.shapes.add_picture(
                avatar_path,
                Inches(6.1), Inches(5.8),
                width=Inches(1.1)
            )

        self.add_fade_transition(slide, 800)

    def slide_2_chart_basic(self):
        """Slide 2: QuickChart.io integration"""
        print("Creating Slide 2: Chart Capabilities...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Chart Generation via QuickChart.io"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Generate chart
        chart_config = {
            "type": "bar",
            "data": {
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "datasets": [{
                    "label": "Revenue (Millions)",
                    "data": [12, 19, 15, 25],
                    "backgroundColor": "rgba(0, 120, 212, 0.7)"
                }]
            },
            "options": {
                "title": {
                    "display": True,
                    "text": "Quarterly Revenue 2025"
                },
                "scales": {
                    "yAxes": [{
                        "ticks": {
                            "beginAtZero": True
                        }
                    }]
                }
            }
        }

        chart_path = f"{self.assets_dir}/chart_q1_q4.png"
        if APIUtils.generate_quickchart(chart_config, chart_path, width=900, height=500):
            slide.shapes.add_picture(
                chart_path,
                Inches(2.5), Inches(2),
                width=Inches(8.3)
            )

        self.add_fade_transition(slide, 600)

    def slide_3_chart_morph(self):
        """Slide 3: Same chart with updated data + morph transition"""
        print("Creating Slide 3: Chart with Morph Transition...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title (same position for morph)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Updated Data (Watch the Morph!)"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Updated chart data
        chart_config = {
            "type": "bar",
            "data": {
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "datasets": [{
                    "label": "Revenue (Millions)",
                    "data": [15, 23, 20, 32],
                    "backgroundColor": "rgba(135, 100, 184, 0.7)"
                }]
            },
            "options": {
                "title": {
                    "display": True,
                    "text": "Quarterly Revenue 2026 (Projected)"
                },
                "scales": {
                    "yAxes": [{
                        "ticks": {
                            "beginAtZero": True
                        }
                    }]
                }
            }
        }

        chart_path = f"{self.assets_dir}/chart_updated.png"
        if APIUtils.generate_quickchart(chart_config, chart_path, width=900, height=500):
            slide.shapes.add_picture(
                chart_path,
                Inches(2.5), Inches(2),
                width=Inches(8.3)
            )

        # Add morph transition
        self.add_morph_transition(slide, 1500)

    def slide_4_ai_images(self):
        """Slide 4: AI-generated images via Ideogram"""
        print("Creating Slide 4: AI-Generated Images...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "AI-Generated Visuals (Ideogram API)"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Generate AI image
        prompt = "Modern minimalist office workspace with laptop, coffee, and natural lighting, professional photography, clean aesthetic"
        image_path = f"{self.assets_dir}/ai_generated_workspace.png"

        if APIUtils.generate_ideogram_image(prompt, image_path, aspect_ratio="ASPECT_16_9"):
            slide.shapes.add_picture(
                image_path,
                Inches(1.5), Inches(2),
                width=Inches(10.3)
            )

            # Caption
            caption_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.6)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = f'Prompt: "{prompt}"'
            for run in caption_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.italic = True
                run.font.color.rgb = self.colors['medium_gray']

        self.add_fade_transition(slide, 600)

    def slide_5_stock_photos(self):
        """Slide 5: Stock photography integration"""
        print("Creating Slide 5: Stock Photography...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Stock Photography (Unsplash + Pexels)"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Unsplash image
        unsplash_path = f"{self.assets_dir}/stock_unsplash.png"
        if APIUtils.get_unsplash_image("teamwork business", unsplash_path, orientation="landscape"):
            slide.shapes.add_picture(
                unsplash_path,
                Inches(0.8), Inches(1.8),
                width=Inches(5.8)
            )

            # Label
            label1 = slide.shapes.add_textbox(
                Inches(0.8), Inches(5.7), Inches(5.8), Inches(0.5)
            )
            label1.text_frame.text = "Unsplash API"
            for run in label1.text_frame.paragraphs[0].runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = self.colors['primary_blue']

        # Pexels image
        pexels_path = f"{self.assets_dir}/stock_pexels.png"
        if APIUtils.get_pexels_image("technology innovation", pexels_path, orientation="landscape"):
            slide.shapes.add_picture(
                pexels_path,
                Inches(6.8), Inches(1.8),
                width=Inches(5.8)
            )

            # Label
            label2 = slide.shapes.add_textbox(
                Inches(6.8), Inches(5.7), Inches(5.8), Inches(0.5)
            )
            label2.text_frame.text = "Pexels API"
            for run in label2.text_frame.paragraphs[0].runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = self.colors['accent_purple']

        self.add_fade_transition(slide, 600)

    def slide_6_diagrams(self):
        """Slide 6: Diagram generation via Mermaid"""
        print("Creating Slide 6: Diagram Generation...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Diagram Generation (Mermaid)"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Mermaid diagram
        mermaid_code = """
        graph TD
            A[User Request] --> B{Presentation Generator}
            B --> C[Generate Outline]
            B --> D[Fetch Visual Assets]
            C --> E[Create Slides]
            D --> E
            E --> F[Apply Transitions]
            F --> G[Export PPTX]

            style A fill:#0078d4,stroke:#0078d4,color:#fff
            style G fill:#107c10,stroke:#107c10,color:#fff
            style B fill:#8764b8,stroke:#8764b8,color:#fff
        """

        diagram_path = f"{self.assets_dir}/mermaid_flowchart.png"
        if APIUtils.generate_mermaid_diagram(mermaid_code, diagram_path, theme="default"):
            slide.shapes.add_picture(
                diagram_path,
                Inches(2), Inches(1.8),
                width=Inches(9.3)
            )

        self.add_fade_transition(slide, 600)

    def slide_7_avatar_grid(self):
        """Slide 7: DiceBear avatar grid"""
        print("Creating Slide 7: Avatar Grid...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Avatar Generation (DiceBear API)"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = self.colors['dark_gray']

        # Generate avatar grid
        avatar_data = [
            ("Alice", "avataaars"),
            ("Bob", "bottts"),
            ("Charlie", "pixel-art"),
            ("Diana", "fun-emoji"),
            ("Eve", "lorelei"),
            ("Frank", "personas"),
        ]

        start_x = 1.5
        start_y = 2.2
        spacing = 2.0
        avatar_size = 1.3

        for i, (name, style) in enumerate(avatar_data):
            col = i % 3
            row = i // 3

            x = start_x + (col * spacing)
            y = start_y + (row * spacing)

            # Generate avatar
            avatar_path = f"{self.assets_dir}/avatar_{name.lower()}.png"
            if APIUtils.get_dicebear_avatar(name, avatar_path, style=style, size=200):
                slide.shapes.add_picture(
                    avatar_path,
                    Inches(x), Inches(y),
                    width=Inches(avatar_size)
                )

                # Name label
                label = slide.shapes.add_textbox(
                    Inches(x), Inches(y + avatar_size + 0.1),
                    Inches(avatar_size), Inches(0.4)
                )
                label.text_frame.text = name
                label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                for run in label.text_frame.paragraphs[0].runs:
                    run.font.size = Pt(14)
                    run.font.bold = True
                    run.font.color.rgb = self.colors['dark_gray']

        self.add_fade_transition(slide, 600)

    def slide_8_full_integration(self):
        """Slide 8: All capabilities combined"""
        print("Creating Slide 8: Full Integration...")

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Title with gradient background shape
        title_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.colors['primary_blue']
        title_bg.line.color.rgb = self.colors['primary_blue']

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Full Integration: All Capabilities Combined"
        for run in title_frame.paragraphs[0].runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = self.colors['white']

        # AI-generated hero image
        hero_prompt = "Abstract technology network connections, data visualization, modern digital art, blue and purple tones"
        hero_path = f"{self.assets_dir}/hero_integration.png"
        if APIUtils.generate_ideogram_image(hero_prompt, hero_path, aspect_ratio="ASPECT_16_9"):
            slide.shapes.add_picture(
                hero_path,
                Inches(0.5), Inches(1.5),
                width=Inches(7)
            )

        # Chart overlay
        mini_chart_config = {
            "type": "line",
            "data": {
                "labels": ["Jan", "Feb", "Mar", "Apr", "May"],
                "datasets": [{
                    "label": "Growth",
                    "data": [10, 15, 13, 20, 25],
                    "borderColor": "rgb(135, 100, 184)",
                    "backgroundColor": "rgba(135, 100, 184, 0.1)"
                }]
            }
        }
        mini_chart_path = f"{self.assets_dir}/mini_chart.png"
        if APIUtils.generate_quickchart(mini_chart_config, mini_chart_path, width=500, height=300):
            slide.shapes.add_picture(
                mini_chart_path,
                Inches(8), Inches(1.5),
                width=Inches(4.5)
            )

        # Bullet points with icons
        bullet_y = 4.8
        bullets = [
            "AI-Generated Images (Ideogram)",
            "Dynamic Charts (QuickChart.io)",
            "Stock Photography (Unsplash/Pexels)",
            "Diagrams & Flowcharts (Mermaid)",
            "Custom Avatars (DiceBear)",
        ]

        for i, bullet_text in enumerate(bullets):
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(bullet_y + i * 0.35),
                Inches(11), Inches(0.3)
            )
            bullet_frame = bullet_box.text_frame
            bullet_frame.text = f"• {bullet_text}"
            for run in bullet_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.color.rgb = self.colors['dark_gray']

        self.add_fade_transition(slide, 600)

    def generate(self, output_path="output/capabilities-showcase.pptx"):
        """Generate the complete showcase presentation"""
        print("=" * 60)
        print("GENERATING CAPABILITIES SHOWCASE PRESENTATION")
        print("=" * 60)

        try:
            self.slide_1_title()
            self.slide_2_chart_basic()
            self.slide_3_chart_morph()
            self.slide_4_ai_images()
            self.slide_5_stock_photos()
            self.slide_6_diagrams()
            self.slide_7_avatar_grid()
            self.slide_8_full_integration()

            # Save presentation
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            self.prs.save(output_path)

            print("\n" + "=" * 60)
            print(f"SUCCESS! Presentation saved to: {output_path}")
            print("=" * 60)
            print("\nSlides created:")
            print("  1. Title Slide with DiceBear Avatar")
            print("  2. QuickChart.io Bar Chart")
            print("  3. Updated Chart with Morph Transition")
            print("  4. AI-Generated Image (Ideogram)")
            print("  5. Stock Photography (Unsplash + Pexels)")
            print("  6. Mermaid Flowchart Diagram")
            print("  7. DiceBear Avatar Grid (6 styles)")
            print("  8. Full Integration Slide")
            print("\nOpen the file in PowerPoint to view transitions!")

            return True

        except Exception as e:
            print(f"\nERROR: Failed to generate presentation")
            print(f"Error details: {e}")
            import traceback
            traceback.print_exc()
            return False


if __name__ == "__main__":
    generator = PresentationGenerator()
    success = generator.generate()
    sys.exit(0 if success else 1)
