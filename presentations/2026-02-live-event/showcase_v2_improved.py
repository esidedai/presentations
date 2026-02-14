"""
V2 Showcase - IMPROVED with eSided Brand Design
Incorporating feedback: logo, website info, curved shapes, no gray backgrounds, readable charts
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent.parent / ".ai" / "showcase"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from lxml import etree
from api_utils import APIUtils

class ShowcaseV2Improved:
    """Showcase demonstrating professional presentation generation with eSided branding"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Paths
        self.project_root = Path(__file__).parent.parent.parent
        self.assets_dir = self.project_root / "assets"
        self.assets_local = Path(__file__).parent / "assets-local"
        self.output_dir = Path(__file__).parent / "output"
        self.output_dir.mkdir(exist_ok=True)

        # Brand colors
        self.brand_blue = RGBColor(0, 120, 212)
        self.brand_purple = RGBColor(135, 100, 184)
        self.dark_gray = RGBColor(50, 49, 48)
        self.light_gray = RGBColor(243, 242, 241)
        self.success_green = RGBColor(16, 124, 16)
        self.white = RGBColor(255, 255, 255)

        # Logo path
        self.logo_path = self.assets_local / "eSided-Logo.png"

    def add_brand_elements(self, slide):
        """Add eSided logo, curved shapes, and website info to slide"""

        # 1. Add curved shape decorations
        # Top curve (blue)
        top_curve = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(-0.5),
            Inches(4), Inches(2)
        )
        top_curve.fill.solid()
        top_curve.fill.fore_color.rgb = self.brand_blue
        top_curve.line.fill.background()
        top_curve.shadow.inherit = False

        # Left curve (purple)
        left_curve = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(3),
            Inches(2), Inches(4)
        )
        left_curve.fill.solid()
        left_curve.fill.fore_color.rgb = self.brand_purple
        left_curve.line.fill.background()
        left_curve.shadow.inherit = False

        # Bottom curve (blue)
        bottom_curve = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(13), Inches(7.5),
            Inches(4), Inches(3)
        )
        bottom_curve.fill.solid()
        bottom_curve.fill.fore_color.rgb = self.brand_blue
        bottom_curve.line.fill.background()
        bottom_curve.shadow.inherit = False

        # 2. Add eSided logo at top
        if self.logo_path.exists():
            logo = slide.shapes.add_picture(
                str(self.logo_path),
                Inches(0.6), Inches(0.35),
                width=Inches(1.5)
            )

        # 3. Add website info at bottom right
        website_box = slide.shapes.add_textbox(
            Inches(11), Inches(8.2),
            Inches(4.5), Inches(0.6)
        )
        website_frame = website_box.text_frame
        website_frame.text = "www.eSided.com | info@esided.com"

        para = website_frame.paragraphs[0]
        para.font.name = 'Inter'
        para.font.size = Pt(12)
        para.font.color.rgb = self.dark_gray
        para.alignment = PP_ALIGN.RIGHT

    def create_table_comparison(self, slide, left_header, right_header, rows):
        """Create a professional 2-column comparison table"""
        # Table dimensions
        x, y = Inches(2.5), Inches(2.5)
        width = Inches(11)
        height = Inches(4.5)

        # Create table
        table = slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=2,
            left=x, top=y,
            width=width, height=height
        ).table

        # Set column widths
        table.columns[0].width = Inches(5.5)
        table.columns[1].width = Inches(5.5)

        # Header row
        header_cells = [table.cell(0, 0), table.cell(0, 1)]
        headers = [left_header, right_header]

        for cell, text in zip(header_cells, headers):
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.dark_gray

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(24)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.white
            paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Data rows
        for i, (left, right) in enumerate(rows, 1):
            # Left cell
            left_cell = table.cell(i, 0)
            left_cell.text = left
            left_cell.fill.solid()
            left_cell.fill.fore_color.rgb = RGBColor(255, 245, 245)  # Light red

            left_para = left_cell.text_frame.paragraphs[0]
            left_para.font.size = Pt(20)
            left_para.font.color.rgb = self.dark_gray
            left_para.alignment = PP_ALIGN.LEFT
            left_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            left_cell.margin_left = Inches(0.3)
            left_cell.margin_right = Inches(0.3)

            # Right cell
            right_cell = table.cell(i, 1)
            right_cell.text = right
            right_cell.fill.solid()
            right_cell.fill.fore_color.rgb = RGBColor(240, 255, 240)  # Light green

            right_para = right_cell.text_frame.paragraphs[0]
            right_para.font.size = Pt(20)
            right_para.font.color.rgb = self.dark_gray
            right_para.alignment = PP_ALIGN.LEFT
            right_cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            right_cell.margin_left = Inches(0.3)
            right_cell.margin_right = Inches(0.3)

    def create_question_slide(self, title, subtext=None):
        """Create a provocative question slide with brand elements"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background (NO GRAY)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

        # Add brand elements
        self.add_brand_elements(slide)

        # Large provocative question
        question_box = slide.shapes.add_textbox(
            Inches(2), Inches(2.5), Inches(12), Inches(3)
        )
        question_frame = question_box.text_frame
        question_frame.text = title
        question_frame.word_wrap = True

        para = question_frame.paragraphs[0]
        para.font.name = 'Playfair Display'
        para.font.size = Pt(54)
        para.font.bold = True
        para.font.color.rgb = self.dark_gray
        para.alignment = PP_ALIGN.CENTER

        # Subtext if provided
        if subtext:
            subtext_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(5.5), Inches(11), Inches(1.5)
            )
            subtext_frame = subtext_box.text_frame
            subtext_frame.text = subtext
            subtext_frame.word_wrap = True

            sub_para = subtext_frame.paragraphs[0]
            sub_para.font.name = 'Inter'
            sub_para.font.size = Pt(26)
            sub_para.font.color.rgb = self.brand_blue
            sub_para.alignment = PP_ALIGN.CENTER

        return slide

    def generate_advanced_chart(self):
        """Generate chart with IMPROVED readability - larger fonts, better spacing"""
        print("Generating improved chart...")

        # Simplified, more readable chart
        chart_config = {
            "type": "bar",
            "data": {
                "labels": ["Platform Cost", "Config & Training", "Value Created", "Net ROI"],
                "datasets": [{
                    "data": [-100, -47, 479.6, 332.6],
                    "backgroundColor": [
                        "rgba(211, 52, 56, 0.85)",
                        "rgba(211, 52, 56, 0.85)",
                        "rgba(16, 124, 16, 0.85)",
                        "rgba(0, 120, 212, 0.95)"
                    ],
                    "borderColor": [
                        "rgb(211, 52, 56)",
                        "rgb(211, 52, 56)",
                        "rgb(16, 124, 16)",
                        "rgb(0, 120, 212)"
                    ],
                    "borderWidth": 3
                }]
            },
            "options": {
                "plugins": {
                    "datalabels": {
                        "anchor": "end",
                        "align": "top",
                        "color": "#000000",
                        "font": {
                            "weight": "bold",
                            "size": 28
                        },
                        "formatter": "(value) => { return value > 0 ? '+$' + value + 'K' : '$' + value + 'K'; }"
                    },
                    "title": {
                        "display": True,
                        "text": "First Year Financial Impact",
                        "font": {
                            "size": 32,
                            "weight": "bold"
                        },
                        "padding": 20
                    },
                    "legend": {
                        "display": False
                    }
                },
                "scales": {
                    "y": {
                        "ticks": {
                            "callback": "(value) => '$' + value + 'K'",
                            "font": {
                                "size": 20
                            },
                            "padding": 10
                        },
                        "grid": {
                            "color": "rgba(0, 0, 0, 0.1)"
                        }
                    },
                    "x": {
                        "ticks": {
                            "font": {
                                "size": 22,
                                "weight": "bold"
                            },
                            "padding": 10
                        },
                        "grid": {
                            "display": False
                        }
                    }
                },
                "layout": {
                    "padding": {
                        "top": 40,
                        "bottom": 20,
                        "left": 20,
                        "right": 20
                    }
                }
            }
        }

        chart_path = self.output_dir / "showcase_chart_improved.png"
        if APIUtils.generate_quickchart(chart_config, str(chart_path), width=1600, height=800):
            print(f"  [OK] Improved chart saved (larger fonts, better spacing)")
            return chart_path
        return None

    def build_showcase(self):
        """Build the IMPROVED V2 showcase presentation"""
        print("="*60)
        print("BUILDING IMPROVED V2 SHOWCASE")
        print("With eSided branding and design improvements")
        print("="*60)

        # SLIDE 1: Professional Title Slide
        print("\n[1/7] Professional title slide with hero image")
        slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        background = slide1.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

        # Add brand elements
        self.add_brand_elements(slide1)

        # Title (left side)
        title_box = slide1.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(7), Inches(2.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Building Your\nAI Workforce"

        for para in title_frame.paragraphs:
            para.font.name = 'Playfair Display'
            para.font.size = Pt(60)
            para.font.bold = True
            para.font.color.rgb = self.dark_gray
            para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 0.9

        # Subtitle
        subtitle_box = slide1.shapes.add_textbox(
            Inches(1), Inches(5.2), Inches(7), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "V2 Showcase: Lessons Learned"

        sub_para = subtitle_frame.paragraphs[0]
        sub_para.font.name = 'Inter'
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = self.brand_blue

        # Hero illustration (right side)
        title_illustration = self.assets_local / "Title-Slide-Illustration.png"
        if title_illustration.exists():
            slide1.shapes.add_picture(
                str(title_illustration),
                Inches(8.5), Inches(1.5),
                width=Inches(7)
            )
            print("  [OK] Using professional title illustration")

        # SLIDE 2: Team Leadership
        print("[2/7] Team slide with actual photos")
        slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        background = slide2.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

        # Add brand elements
        self.add_brand_elements(slide2)

        # Title
        title_box = slide2.shapes.add_textbox(
            Inches(2.5), Inches(0.8), Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Leadership"

        title_para = title_frame.paragraphs[0]
        title_para.font.name = 'Playfair Display'
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray
        title_para.alignment = PP_ALIGN.CENTER

        # Team members (2x2 grid)
        team = [
            {"name": "Zak Ali", "title": "CEO, AI & Transformation Expert", "photo": "Zak-Ali.png"},
            {"name": "Dr. Humaira Waqas", "title": "CFO", "photo": "Dr-Humaira-Waqas.png"},
            {"name": "Dr. Kamran Malik", "title": "CTO", "photo": "Dr-Kamran-Malik.png"},
            {"name": "Olga Skeen", "title": "COO", "photo": "Olga-Skeen.png"}
        ]

        x_start, y_start = 2.5, 2.2
        x_spacing, y_spacing = 3.2, 3.2

        for i, person in enumerate(team):
            row = i // 2
            col = i % 2
            x = x_start + (col * x_spacing)
            y = y_start + (row * y_spacing)

            # Photo
            photo_path = self.assets_local / "Team-Pictures" / person['photo']
            if photo_path.exists():
                slide2.shapes.add_picture(
                    str(photo_path),
                    Inches(x), Inches(y),
                    width=Inches(2.2), height=Inches(2.2)
                )

            # Name
            name_box = slide2.shapes.add_textbox(
                Inches(x), Inches(y + 2.35),
                Inches(2.2), Inches(0.35)
            )
            name_frame = name_box.text_frame
            name_frame.text = person['name']
            name_para = name_frame.paragraphs[0]
            name_para.font.name = 'Inter'
            name_para.font.size = Pt(16)
            name_para.font.bold = True
            name_para.font.color.rgb = self.dark_gray
            name_para.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide2.shapes.add_textbox(
                Inches(x), Inches(y + 2.7),
                Inches(2.2), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = person['title']
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Inter'
            title_para.font.size = Pt(12)
            title_para.font.color.rgb = RGBColor(142, 142, 147)
            title_para.alignment = PP_ALIGN.CENTER

        print("  [OK] Team photos and bios added")

        # SLIDE 3: Provocative Question
        print("[3/7] Provocative question slide")
        self.create_question_slide(
            "What's the equivalent\nin YOUR business\ntoday?",
            "What work are your people doing that should be eliminated?"
        )

        # SLIDE 4: Workers vs Tools (TABLE)
        print("[4/7] Comparison table (Workers vs Tools)")
        slide4 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background
        background = slide4.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

        # Add brand elements
        self.add_brand_elements(slide4)

        # Title
        title_box = slide4.shapes.add_textbox(
            Inches(2.5), Inches(0.8), Inches(11), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Tools vs. Workers"

        title_para = title_frame.paragraphs[0]
        title_para.font.name = 'Playfair Display'
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray
        title_para.alignment = PP_ALIGN.CENTER

        # Comparison table
        self.create_table_comparison(
            slide4,
            "Tools (ChatGPT)",
            "Workers (AI Workforce)",
            [
                ("You start every task", "Runs autonomously on triggers"),
                ("No memory (stateless)", "Remembers patterns & history"),
                ("You provide context every time", "Learns institutional knowledge"),
                ("Output: Suggestions", "Output: Completed work product")
            ]
        )
        print("  [OK] Professional comparison table created")

        # SLIDE 5: Advanced Chart
        print("[5/7] Improved chart with larger fonts and better spacing")
        chart_path = self.generate_advanced_chart()

        if chart_path:
            slide5 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            # White background
            background = slide5.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.white

            # Add brand elements
            self.add_brand_elements(slide5)

            # Title
            title_box = slide5.shapes.add_textbox(
                Inches(2.5), Inches(0.8), Inches(11), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = "The Numbers"

            title_para = title_frame.paragraphs[0]
            title_para.font.name = 'Playfair Display'
            title_para.font.size = Pt(48)
            title_para.font.bold = True
            title_para.font.color.rgb = self.dark_gray
            title_para.alignment = PP_ALIGN.CENTER

            # Chart (centered, larger)
            slide5.shapes.add_picture(
                str(chart_path),
                Inches(2), Inches(2),
                width=Inches(12)
            )

            # Key insight callout
            insight_box = slide5.shapes.add_textbox(
                Inches(2.5), Inches(7.3), Inches(11), Inches(0.7)
            )
            insight_frame = insight_box.text_frame
            insight_frame.text = "3.7 month payback  •  221 patterns encoded  •  Permanent institutional asset"

            insight_para = insight_frame.paragraphs[0]
            insight_para.font.name = 'Inter'
            insight_para.font.size = Pt(18)
            insight_para.font.bold = True
            insight_para.font.color.rgb = self.brand_blue
            insight_para.alignment = PP_ALIGN.CENTER

        # SLIDE 6: Short Provocative Statement
        print("[6/7] Short provocative statement")
        slide6 = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # White background (NO GRAY)
        background = slide6.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.white

        # Add brand elements
        self.add_brand_elements(slide6)

        # Statement
        statement_box = slide6.shapes.add_textbox(
            Inches(2.5), Inches(3), Inches(11), Inches(2.5)
        )
        statement_frame = statement_box.text_frame
        statement_frame.text = "None of us mourn\nthese jobs."
        statement_frame.word_wrap = True

        for para in statement_frame.paragraphs:
            para.font.name = 'Playfair Display'
            para.font.size = Pt(64)
            para.font.bold = True
            para.font.color.rgb = self.dark_gray
            para.alignment = PP_ALIGN.CENTER
            para.line_spacing = 1.0

        # Subtext
        subtext_box = slide6.shapes.add_textbox(
            Inches(2.5), Inches(5.8), Inches(11), Inches(1)
        )
        subtext_frame = subtext_box.text_frame
        subtext_frame.text = "The work was eliminated. The workers weren't."

        sub_para = subtext_frame.paragraphs[0]
        sub_para.font.name = 'Inter'
        sub_para.font.size = Pt(26)
        sub_para.font.color.rgb = self.brand_blue
        sub_para.alignment = PP_ALIGN.CENTER

        # SLIDE 7: Final Question
        print("[7/7] Final provocative question")
        self.create_question_slide(
            "What's your stance\non this transition?",
            "Resisting? Accepting? Participating?"
        )

        # Save presentation
        output_file = self.output_dir / "V2_Showcase_Improved.pptx"
        self.prs.save(str(output_file))

        print("\n" + "="*60)
        print("IMPROVED V2 SHOWCASE COMPLETE")
        print("="*60)
        print(f"Saved to: {output_file}")
        print(f"Total slides: {len(self.prs.slides)}")
        print("\nDesign improvements:")
        print("  [OK] eSided logo on all slides (top)")
        print("  [OK] Website info on all slides (bottom right)")
        print("  [OK] Curved decorative shapes (blue/purple)")
        print("  [OK] White backgrounds (NO GRAY)")
        print("  [OK] Improved chart readability (larger fonts, better spacing)")
        print("  [OK] Professional typography (Playfair + Inter)")
        print("\nReady for feedback!")

        return output_file


def main():
    showcase = ShowcaseV2Improved()
    showcase.build_showcase()


if __name__ == "__main__":
    main()
