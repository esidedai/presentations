"""
PNG to PowerPoint Assembler
Inserts rendered PNGs as full-slide images with speaker notes for rehearsal.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


# Enhanced speaker notes for Presenter View rehearsal.
# Format: WHY this slide matters → WHAT to say → HOW to deliver it.
# Transition cues marked with >>. Key phrases in CAPS for emphasis.
SPEAKER_NOTES = {
    "01-title": (
        "WHY: Set the frame immediately — this talk is NOT about AI tools. "
        "The audience has heard a dozen 'AI will change everything' talks. "
        "You need to differentiate in the first 30 seconds.\n\n"
        "SAY: 'Thanks for being here. I want to start with something that "
        "might seem unrelated to your business. Bear with me.'\n\n"
        "DELIVER: Confident, conversational. Make eye contact across the room. "
        "Don't linger on the title slide — move quickly to the pattern."
    ),
    "02-combine-harvester": (
        "WHY: You're opening with history to build a PATTERN the audience "
        "will recognize later. This isn't nostalgia — it's setup for the "
        "punchline that comes on slides 5-6.\n\n"
        "SAY: 'In 1900, harvesting a hundred acres took dozens of laborers. "
        "Weeks of backbreaking work. Children in fields with scythes. "
        "Today: one combine, one operator, air-conditioned cab, GPS-guided. "
        "That transition eliminated work. But nobody wishes we still had "
        "children in fields with scythes.'\n\n"
        "DELIVER: Slow, storytelling pace. Let the image breathe. "
        ">> TRANSITION: 'Same pattern, different era...'"
    ),
    "03-switchboard-operators": (
        "WHY: Switchboard operators are the KNOWLEDGE WORK example. "
        "Agriculture was physical labor — easy for audience to dismiss. "
        "Switchboard operators used skill, judgment, memory — just like "
        "their senior staff. This makes it personal.\n\n"
        "SAY: 'You picked up the phone and a voice answered: \"Number please.\" "
        "200,000 operators in 1920. Entire buildings full of them. "
        "Then came direct dial. By the 1980s, the job was gone. "
        "Nobody wants to go back.'\n\n"
        "DELIVER: The quote 'Number please' should feel like a callback. "
        "Pause after 'Nobody wants to go back.' Let them sit with it."
    ),
    "04-human-computers": (
        "WHY: This one hits closest to home — 'Computer' as a job title "
        "is the most visceral example. It connects to the audience's own "
        "professional identity. It's also the closing callback.\n\n"
        "SAY: 'Here's one that might hit closer to home. Before the electronic "
        "computer existed, \"computer\" was a job title. Capital C. Printed "
        "on civil service paperwork. 30 to 40 hours of hand calculations "
        "for a single trajectory. Now your phone does that in milliseconds. "
        "The machine has the name.'\n\n"
        "DELIVER: Emphasize 'Capital C.' and 'job title' — these are the "
        "phrases that come back on the closing slide. Plant them now."
    ),
    "05-the-landing": (
        "WHY: This is where you land the pattern and pivot to the present. "
        "Three historical examples = undeniable pattern. Now apply it to "
        "the audience. The emotional weight is in 'deserved better.'\n\n"
        "SAY: 'None of us mourn these jobs. The work that was eliminated was "
        "work that should have been eliminated. The people deserved better "
        "than that work. Their grandchildren are accountants, engineers, "
        "attorneys. Some of them are probably in this room. "
        "AI is the next one.'\n\n"
        "DELIVER: 'AI is the next one' is the bridge. Say it with certainty, "
        "not drama. Then click to the next slide for the challenge."
    ),
    "06-ai-is-next": (
        "WHY: Turning the mirror on the audience. The historical examples "
        "were safe — now it's about THEIR business. This creates productive "
        "tension that carries the rest of the talk.\n\n"
        "SAY: 'What's the equivalent in YOUR business today? What work should "
        "be eliminated? What work are your people doing that they deserve "
        "better than?'\n\n"
        "DELIVER: Make eye contact with specific people. These are genuine "
        "questions — let them land. Pause 3 seconds before moving on."
    ),
    "07-pre-qualification": (
        "WHY: Setting expectations — this is NOT a vendor pitch. The audience "
        "needs permission to think strategically instead of defensively. "
        "This builds trust and lowers resistance.\n\n"
        "SAY: 'This session is for business owners and firm leaders who want "
        "to understand where AI actually fits, and where it doesn't. "
        "If you're looking for a vendor pitch, this isn't that. "
        "We're going to talk about how to think about this.'\n\n"
        "DELIVER: Conversational, reassuring. 'This isn't that' with a "
        "slight head shake. Lower the sales guard."
    ),
    "08-participate-well": (
        "WHY: Introducing the recurring motif that closes the talk. "
        "'Participate well' is the phrase they should remember. "
        "You're also naming their fears (outcompeted, people impact, "
        "past failures) before they voice them.\n\n"
        "SAY: 'You're inside a transition that has happened before. "
        "You know how these end. Your job isn't to resist or accept. "
        "Your job is to participate well. So how do you participate well? "
        "That's what the rest of this is about. But first, there's "
        "something in the way.'\n\n"
        "DELIVER: 'Participate well' — say it slowly, with emphasis. "
        "This phrase comes back on slide 34. Plant it now. "
        ">> TRANSITION: 'But first, there's something in the way...'"
    ),
    "09-technology-before": (
        "WHY: Meeting the audience where they are. They HAVE a mental model "
        "for technology. You're naming it explicitly so you can challenge "
        "it. This builds credibility — you understand their world.\n\n"
        "SAY: 'You've bought technology before. You have a mental model. "
        "Technology is a tool. Evaluate features. Compare vendors. "
        "Buy, configure, roll out. This frame has been sold to you "
        "for decades.'\n\n"
        "DELIVER: List the steps casually — they'll nod. 'Sold to you' "
        "plants the seed that this frame might be wrong."
    ),
    "10-crm-promise": (
        "WHY: CRM is the perfect analogy because EVERYONE in the room has "
        "bought CRM. It's familiar. The promise was real. The outcome "
        "varied. This setup makes the AI pivot land harder.\n\n"
        "SAY: 'CRM promised revenue growth. \"Buy this, your revenue will "
        "grow.\" The technology causes the outcome. Salesforce is worth "
        "$200 billion. It works. But what does it actually deliver? "
        "Visibility. Reporting. Data structure. That's real value. "
        "But it's infrastructure.'\n\n"
        "DELIVER: The quote should sound like a sales pitch. "
        "Then undercut it gently with 'But it's infrastructure.'"
    ),
    "11-what-actually-happened": (
        "WHY: Revealing the truth behind the CRM promise. Same tech, "
        "different outcomes. The audience knows this from experience — "
        "you're giving them language for something they felt but "
        "couldn't articulate.\n\n"
        "SAY: 'The companies that grew had something else underneath. "
        "A sales process that worked. Management that paid attention. "
        "The CRM amplified what was already there. Same technology. "
        "Wildly different outcomes. The difference was invisible "
        "at purchase time.'\n\n"
        "DELIVER: 'Invisible at purchase time' is the key insight. "
        "Pause after it. Some will lean forward."
    ),
    "12-crm-vs-people": (
        "WHY: THIS IS THE EMOTIONAL PIVOT of the entire talk. The left "
        "column is what systems capture. The right column is where the "
        "real value lives. Every person in the room has a 'Sarah.' "
        "This slide creates the problem that AI workers solve.\n\n"
        "SAY: [Point to left column] 'This is what HubSpot holds. "
        "Contact name. Deal value: $60,000. Status: Proposal Sent.' "
        "[Point to right column] 'This is what your BD director actually "
        "knows. The CEO just went through a leadership transition. "
        "This deal is really worth $180K. The CFO takes three weeks "
        "to review anything. Following up now would be counterproductive.' "
        "[Pause] 'The right column is where the actual value lives. "
        "No system built around database fields can capture it.'\n\n"
        "DELIVER: SLOW. Let the contrast do the work. 'Where does that "
        "knowledge live when she leaves?' should be almost quiet. "
        "This is the gut punch of the talk."
    ),
    "13-amplifier": (
        "WHY: Naming the principle. Technology amplifies what's there. "
        "Good process + tech = better. No process + tech = expensive "
        "nothing. This is familiar — they've lived it.\n\n"
        "SAY: 'Technology is an amplifier. It amplifies what's there. "
        "Good process plus CRM equals better process. No process plus "
        "CRM equals expensive nothing. And the right column? Sarah's "
        "patterns, Tom's vendor knowledge, Katie's hiring instincts? "
        "That walks out the door when they leave.'\n\n"
        "DELIVER: 'Walks out the door' — some will wince. They know "
        "this pain personally. Use it."
    ),
    "14-ai-arrives": (
        "WHY: Setting up the trap. Same pattern loading = same mental "
        "model being applied. 'But there's a compounding problem' "
        "creates anticipation for the reveal.\n\n"
        "SAY: 'Now AI arrives. Same pattern loading. \"Buy AI. Automate "
        "workflows. Get productivity.\" But there's a compounding problem.'\n\n"
        "DELIVER: Quick delivery. The quote should sound familiar — like "
        "a vendor pitch they've heard. 'Compounding problem' with a "
        "slight lean forward. >> TRANSITION: Click to the reveal."
    ),
    "15-wrong-category": (
        "WHY: The intellectual core of the argument. CRM = tool (correct). "
        "AI = tool (WRONG). The four traits (initiates, judgment, autonomy, "
        "develops) are each a 'wait, that's different' moment.\n\n"
        "SAY: 'For CRM, \"tool\" was the right category. The expectation "
        "was wrong. For AI, \"tool\" is the wrong category entirely. "
        "AI initiates — it doesn't wait for you to act. AI has judgment. "
        "AI operates with autonomy. AI develops over time. "
        "This doesn't describe a tool. It describes something else.'\n\n"
        "DELIVER: Read each trait distinctly. Let each one register. "
        "'Something else' with a pause — then click to the answer."
    ),
    "16-ai-workers": (
        "WHY: The REFRAME moment. This is where the thesis crystallizes. "
        "Not software. Workers. The three questions (job, supervision, "
        "workflow) are the new mental model.\n\n"
        "SAY: 'AI workers are workers. Not software. Workers. "
        "You don't ask \"what can this tool do?\" You ask: What's the job? "
        "Who supervises? Is the workflow ready?'\n\n"
        "DELIVER: 'Not software. Workers.' — firm, definitive. "
        "The three questions should feel like a checklist they can use. "
        "This is actionable, not theoretical."
    ),
    "17-worker-gets-right": (
        "WHY: Making the worker frame concrete with real examples. "
        "Invoice processing, pipeline building, talent hunting — these "
        "are roles the audience RECOGNIZES in their own firms.\n\n"
        "SAY: 'One outcome per worker. Invoice processing worker owns: "
        "invoices processed accurately and timely. Pipeline builder owns: "
        "qualified prospects surfaced. Talent hunter owns: qualified "
        "candidates surfaced for open roles.'\n\n"
        "DELIVER: Read each outcome like a job description. "
        "'Owns' is the key verb — workers OWN outcomes, tools don't."
    ),
    "18-different-questions": (
        "WHY: Side-by-side comparison makes the frame shift visible. "
        "Tool frame = one question. Worker frame = three questions. "
        "The worker frame is richer, more specific, more actionable.\n\n"
        "SAY: 'Tool frame asks: Does this software have invoice processing? "
        "Worker frame asks: What outcome does this worker own? "
        "What decisions can they make alone? Where do they need supervision? "
        "Different questions. Different downstream.'\n\n"
        "DELIVER: Gesture left for tool frame, right for worker frame. "
        "'Different downstream' should feel like a revelation."
    ),
    "19-you-already-know": (
        "WHY: Giving the audience permission. They ALREADY know management. "
        "They don't need to learn a new framework — just translate the "
        "one they have. This removes the 'it's too complicated' objection.\n\n"
        "SAY: 'The framework is management. Define the role. New hires "
        "ramp up. Some need supervision. Others run. You have the "
        "mental model. It needs translating, not learning.'\n\n"
        "DELIVER: Reassuring tone. 'Translating, not learning' is the "
        "key phrase. They should feel relief, not overwhelm."
    ),
    "20-autonomy-as-design": (
        "WHY: Addressing the unspoken fear — 'Will AI go rogue?' No. "
        "YOU calibrate the autonomy. This slide gives the audience "
        "CONTROL, which is what skeptics need to hear.\n\n"
        "SAY: 'Some workers you want autonomous. They run, you see "
        "exceptions. Some supervised — they recommend, you approve. "
        "Some human-led — they assist, they don't own. "
        "Same underlying technology. Different relationship. Your choice "
        "based on what the work requires.'\n\n"
        "DELIVER: Emphasize 'YOUR choice.' This is about control, "
        "not surrender. Walk through each mode briefly."
    ),
    "21-value-spectrum": (
        "WHY: Introducing the 7-layer framework that most AI conversations "
        "miss. This visual shows WHERE AI operates and WHERE humans stay. "
        "It's not replacement — it's elevation.\n\n"
        "SAY: 'Everything your people carry falls on a spectrum. "
        "Seven layers. AI handles the bottom three: data, information, "
        "knowledge. The middle — expertise and hindsight — AI and humans "
        "share. The top — insight and foresight — that stays human. "
        "What does this pattern mean? What should we do about it?'\n\n"
        "DELIVER: Point to the spectrum visual. Don't lecture all seven. "
        "The visual is the point. Emphasize 'stays human.'"
    ),
    "22-gradient-shifts": (
        "WHY: The boundary MOVES over time. This is the compounding "
        "argument. Week 1 is modest. By year 1, humans focus on "
        "the highest-level work. The system gets better.\n\n"
        "SAY: 'The boundary moves. Week one, AI handles data and "
        "information. By month three, most knowledge work too. "
        "By month six, routine expertise patterns. Your team doesn't "
        "do less work. They do higher-level work.'\n\n"
        "DELIVER: Trace the progression left to right. "
        "'Higher-level work' with conviction."
    ),
    "23-integration-layer": (
        "WHY: Naming the current reality that everyone feels but nobody "
        "talks about. Your people are GLUE between disconnected systems. "
        "That's layers 1-3 work. It should be automated.\n\n"
        "SAY: 'Right now, your people are the integration layer. "
        "Entering data in five different systems. Exporting, importing. "
        "None of them talk to each other. Your people are the glue. "
        "That's work the machines should do.'\n\n"
        "DELIVER: List the tasks with slight frustration — mirror what "
        "the audience feels. 'Machines should do' firmly."
    ),
    "24-judgment-layer": (
        "WHY: The visual shift from constrained (integration) to elevated "
        "(judgment). This is the aspirational promise. The arrow from "
        "one to the other IS the value proposition.\n\n"
        "SAY: 'Your people become the judgment layer. See exceptions. "
        "Make decisions. Do work only humans should do. "
        "Integration layer to judgment layer. That's the shift. "
        "Let me show you what this looks like.'\n\n"
        "DELIVER: 'That's the shift' with finality. "
        "'Let me show you' is the bridge to the demo. "
        ">> TRANSITION: Click to demo transition slide."
    ),
    "25-demo-transition": (
        "WHY: Setting up the demo with context. The audience needs to "
        "know WHO Meridian is and WHY they should care. '50-person firm' "
        "feels relatable. 'Monday morning' creates immediacy.\n\n"
        "SAY: 'Meridian Group. 50-person professional services firm. "
        "$10.5 million in revenue. Five AI workers running real workflows. "
        "You're the managing partner. Three people report to you. "
        "Tom, Sarah, Katie. Every worker you're about to see is based "
        "on a production system. The data is synthetic. "
        "The architecture is not. It's Monday morning. "
        "This is what's waiting for you.'\n\n"
        "DELIVER: Build energy. 'Monday morning' with a grin. "
        "Then switch to the demo. [SWITCH TO DEMO — 9-12 minutes]"
    ),
    "26-three-views": (
        "WHY: Bridge from demo back to slides. Quick recap of the three "
        "tabs they just saw. Don't linger — the demo showed it. "
        "This is a visual anchor for what comes next.\n\n"
        "SAY: 'You just saw three ways to measure what this workforce "
        "is worth. The financial case: $332,600 a year. The institutional "
        "knowledge case: 221 patterns encoded. And the model for your "
        "firm. Different numbers, same methodology.'\n\n"
        "DELIVER: Point to each card. Quick, confident. "
        ">> TRANSITION: 'Here's what most ROI conversations miss...'"
    ),
    "27-spreadsheet-misses": (
        "WHY: Reinforcing the Capital View for anyone who didn't fully "
        "absorb it during the demo. The number 221 should be prominent. "
        "This is the 'appreciates, not depreciates' argument.\n\n"
        "SAY: 'The $332K pays for the investment. That's important. "
        "But it's not why you stay. You stay because 221 patterns are "
        "now organizational knowledge. Tom's 12 years — 52 patterns. "
        "Sarah's deal judgment — 109 patterns. Most technology "
        "depreciates. This does the opposite. It appreciates.'\n\n"
        "DELIVER: 'Appreciates' is the money word. Let it land. "
        "Every time someone teaches the system, it gets smarter."
    ),
    "28-scar-tissue-test": (
        "WHY: The emotional closer of the ROI section. Everyone in the "
        "room has had a key person leave. They know the crater. "
        "This slide shows what fills it.\n\n"
        "SAY: 'What happens when Tom leaves after 12 years? Without AI "
        "workers: six months relearning. Exception rate resets to 15%. "
        "With AI workers: 52 patterns preserved. Starts at month 6. "
        "Exception rate stays at 4%. Same for Sarah. Same for Katie. "
        "The knowledge stays even when the people move on.'\n\n"
        "DELIVER: Read the 'without' column with weight — they've felt "
        "this. Read the 'with' column with relief. 'Knowledge stays' "
        "should feel like a promise."
    ),
    "29-three-questions": (
        "WHY: Prescriptive takeaway #1. Three simple questions before any "
        "AI adoption. Gives the audience an ACTION ITEM they can use "
        "Monday morning. Practical, not theoretical.\n\n"
        "SAY: 'Before you adopt AI for any workflow, three questions. "
        "Is this process documented and measurable? If you can't describe "
        "it, you can't delegate it. What outcome does this worker own? "
        "One sentence. Clear accountability. Who supervises, and when do "
        "they intervene? Autonomy is designed, not assumed.'\n\n"
        "DELIVER: Count each on your fingers. These should feel like "
        "a checklist they can take home. 'Designed, not assumed' lands."
    ),
    "30-three-skills": (
        "WHY: Prescriptive takeaway #2. What their PEOPLE need to develop. "
        "This answers 'what do I tell my team?' Articulation is the most "
        "valuable — every explanation is a deposit.\n\n"
        "SAY: 'For your people, three skills. Articulation — explain your "
        "decisions so the system learns. \"I rejected this lead because PE "
        "rollups under 18 months always churn.\" Every explanation is a "
        "deposit. Evaluation — judge AI output rather than producing from "
        "scratch. Direction-setting — spend time on insight and foresight.'\n\n"
        "DELIVER: The quote is the best example of articulation. "
        "Read it like someone actually explaining their reasoning. "
        "'Every explanation is a deposit' — repeat if needed."
    ),
    "31-deeper-call": (
        "WHY: Callbacks to the opening. Combine + switchboard + computers. "
        "Now applied to THEIR Tom, THEIR Sarah, THEIR Katie. "
        "This is the emotional peak before the honesty slide.\n\n"
        "SAY: 'Your Tom is spending four hours Monday matching invoices. "
        "Your Sarah is spending evenings researching prospects. "
        "Your Katie is waiting 73 days to fill a role. "
        "Your people are the integration layer right now. "
        "They deserve to be the judgment layer.'\n\n"
        "DELIVER: 'Your Tom, Your Sarah, Your Katie' — make it personal. "
        "'They deserve' with quiet conviction. This is the moral argument."
    ),
    "32-honesty": (
        "WHY: Direct challenge. Are you resisting, accepting, or "
        "participating? This requires self-reflection. It's uncomfortable "
        "by design — you want them thinking, not comfortable.\n\n"
        "SAY: 'This requires honesty. What's your stance on this "
        "transition? Are you resisting? Accepting? Participating? "
        "If you're honest with yourself, you can be honest with your team.'\n\n"
        "DELIVER: Slow. Direct eye contact. Let the silence work. "
        "Don't rush past the discomfort."
    ),
    "33-the-close": (
        "WHY: The callback that ties everything together. 'Computer' as a "
        "job title was planted on slide 4. Now it returns as the closing "
        "metaphor. 'You're in the middle' makes it present tense.\n\n"
        "SAY: 'In 1942, \"Computer\" was a job title. Capital C. "
        "Printed on civil service paperwork. 30 to 40 hours of hand "
        "calculations. Now your phone does that in milliseconds. "
        "And nobody mourns the job. You're in the middle of that "
        "transition.'\n\n"
        "DELIVER: SLOW. Each sentence gets its own beat. "
        "'You're in the middle' should land like a mirror."
    ),
    "34-participate-well": (
        "WHY: The closing motif. 'Participate well' was introduced on "
        "slide 8 and now it returns as the final call to action. "
        "Simple, memorable, actionable.\n\n"
        "SAY: 'Your job isn't to resist or accept. "
        "Your job is to participate well.'\n\n"
        "DELIVER: These two sentences. That's it. Don't add. "
        "Don't explain. Let the silence hold for 3-5 seconds. "
        "Then click to Q&A."
    ),
    "35-qa": (
        "WHY: Opening the floor. The talk lands differently for different "
        "people. Q&A is where the real connections happen.\n\n"
        "SAY: 'I'd love to hear your questions.'\n\n"
        "DELIVER: Warm, open. During Q&A, use the demo tools: "
        "info modals for methodology, chat for pre-built answers, "
        "calculator to model audience members' firms live, "
        "Capital View for detailed pattern walkthroughs. "
        "Close with: 'Before any AI adoption — three questions. "
        "Process documented? Outcome clear? Supervision designed? "
        "Your job is to participate well. Thanks for your time.'"
    ),
}


class PowerPointAssembler:
    """Assemble PNG slides into a PowerPoint presentation with speaker notes"""

    def __init__(self, png_dir, output_file):
        self.png_dir = Path(png_dir)
        self.output_file = Path(output_file)
        self.prs = Presentation()
        # Standard 16:9 widescreen
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_slide(self, png_file):
        """Add a PNG as a full-slide image with optional speaker notes"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        slide.shapes.add_picture(
            str(png_file),
            Inches(0), Inches(0),
            width=self.prs.slide_width,
            height=self.prs.slide_height,
        )

        # Add speaker notes if available
        slug = png_file.stem
        if slug in SPEAKER_NOTES:
            slide.notes_slide.notes_text_frame.text = SPEAKER_NOTES[slug]

        return slide

    def assemble_all(self):
        png_files = sorted(self.png_dir.glob("*.png"))
        if not png_files:
            print(f"No PNG files found in {self.png_dir}")
            return False

        print("=" * 60)
        print("POWERPOINT ASSEMBLER")
        print("=" * 60)
        print(f"PNGs:   {self.png_dir}")
        print(f"Output: {self.output_file}")
        print(f"Found {len(png_files)} slides")
        print("=" * 60)

        notes_count = 0
        for png_file in png_files:
            try:
                self.add_slide(png_file)
                has_notes = png_file.stem in SPEAKER_NOTES
                if has_notes:
                    notes_count += 1
                print(f"  [OK] {png_file.name}" + (" + notes" if has_notes else ""))
            except Exception as e:
                print(f"  [X]  {png_file.name}: {e}")

        self.output_file.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(self.output_file))

        print(f"\nCOMPLETE: {len(self.prs.slides)} slides, {notes_count} with speaker notes")
        print(f"Saved to: {self.output_file}")
        return True


def main():
    base = Path(__file__).parent
    assembler = PowerPointAssembler(
        png_dir=base / "output",
        output_file=base.parent / "output" / "AI-Workforce-Deck.pptx",
    )
    success = assembler.assemble_all()
    if success:
        print(f"\n[OK] Open in PowerPoint -> Presenter View to see speaker notes")
    else:
        print("\n[X] Assembly failed")
        return 1
    return 0


if __name__ == "__main__":
    exit(main())
