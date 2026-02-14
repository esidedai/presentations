"""
AI Workforce - Live Event Presentation Generator
Creates a 28-slide presentation for 45-minute live event
"""

import sys
from pathlib import Path

# Add parent directories to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent / ".ai" / "showcase"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree
import os
from api_utils import APIUtils

class AIWorkforcePresentation:
    """Generate the AI Workforce presentation"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

        # Paths
        self.project_root = Path(__file__).parent.parent.parent
        self.assets_dir = self.project_root / "assets"
        self.output_dir = Path(__file__).parent / "output"
        self.output_dir.mkdir(exist_ok=True)

        # Brand colors
        self.brand_blue = RGBColor(0, 120, 212)  # #0078d4
        self.brand_purple = RGBColor(135, 100, 184)  # #8764b8
        self.dark_gray = RGBColor(50, 49, 48)  # #323130
        self.light_gray = RGBColor(243, 242, 241)  # #f3f2f1

    def add_morph_transition(self, slide, duration_ms=1000):
        """Add morph transition to slide"""
        xml = f'''<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
          <mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" Requires="p159">
            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
                          spd="slow" p14:dur="{duration_ms}">
              <p159:morph option="byObject"/>
            </p:transition>
          </mc:Choice>
          <mc:Fallback>
            <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="slow">
              <p:fade/>
            </p:transition>
          </mc:Fallback>
        </mc:AlternateContent>'''
        slide.element.append(etree.fromstring(xml))

    def add_title_slide(self, title, subtitle):
        """Create a title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.light_gray

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(66)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(14), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(32)
            subtitle_para.font.color.rgb = self.brand_blue
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(self, title, content, image_path=None, layout="title_content"):
        """Create a content slide with optional image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray

        # Content area
        if image_path and Path(image_path).exists():
            # Image on right, content on left
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(7), Inches(6)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for line in content:
                p = content_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(24)
                p.font.color.rgb = self.dark_gray
                p.space_before = Pt(12)
                if line.startswith("•"):
                    p.level = 1

            # Add image
            slide.shapes.add_picture(
                str(image_path),
                Inches(8.5), Inches(2),
                width=Inches(7)
            )
        else:
            # Content centered
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5), Inches(13), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for line in content:
                p = content_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(28)
                p.font.color.rgb = self.dark_gray
                p.space_before = Pt(18)
                p.alignment = PP_ALIGN.CENTER
                if line.startswith("•"):
                    p.level = 1
                    p.alignment = PP_ALIGN.LEFT

        return slide

    def add_image_slide(self, title, image_path, caption=None):
        """Create a slide with large image"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray

        # Image
        if Path(image_path).exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(2), Inches(2),
                width=Inches(12)
            )

        # Caption
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(2), Inches(7.5), Inches(12), Inches(1)
            )
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_para = caption_frame.paragraphs[0]
            caption_para.font.size = Pt(28)
            caption_para.font.color.rgb = self.dark_gray
            caption_para.alignment = PP_ALIGN.CENTER

        return slide

    def generate_historical_images(self):
        """Generate historical images using Ideogram"""
        print("\nGenerating historical images...")

        images = {
            "harvest": "Photorealistic historical scene of 1800s harvest with many workers manually harvesting wheat in a field, golden hour lighting, documentary style",
            "combine": "Vintage 1930s combine harvester machine in wheat field, black and white photograph style, industrial documentation",
            "switchboard": "1950s Bell System switchboard operators at work, rows of women operating manual telephone switchboards, vintage photograph style",
            "direct_dial": "Evolution from rotary phone to modern smartphone, clean product photography style, side by side comparison"
        }

        generated = {}
        for name, prompt in images.items():
            output_path = self.output_dir / f"historical_{name}.png"
            print(f"  Generating {name}...")
            try:
                if APIUtils.generate_ideogram_image(prompt, str(output_path)):
                    generated[name] = output_path
                    print(f"    [OK] Saved to {output_path}")
            except Exception as e:
                print(f"    [X] Failed: {e}")
                # Fallback to stock photos
                print(f"    Attempting Unsplash fallback...")
                try:
                    if APIUtils.get_unsplash_image(name, str(output_path)):
                        generated[name] = output_path
                        print(f"    [OK] Saved from Unsplash")
                except:
                    print(f"    [X] Unsplash also failed")

        return generated

    def generate_charts(self):
        """Generate charts using QuickChart"""
        print("\nGenerating charts...")

        charts = {}

        # ROI Breakdown Chart
        roi_config = {
            "type": "bar",
            "data": {
                "labels": ["Platform Cost", "Configuration", "Value Created", "Net ROI"],
                "datasets": [{
                    "label": "Amount ($K)",
                    "data": [-100, -47, 479.6, 332.6],
                    "backgroundColor": [
                        "rgba(255, 99, 132, 0.7)",
                        "rgba(255, 99, 132, 0.7)",
                        "rgba(75, 192, 192, 0.7)",
                        "rgba(54, 162, 235, 0.7)"
                    ]
                }]
            },
            "options": {
                "title": {
                    "display": True,
                    "text": "Meridian Group - First Year ROI",
                    "fontSize": 20
                },
                "scales": {
                    "yAxes": [{
                        "ticks": {
                            "beginAtZero": True,
                            "fontSize": 16
                        }
                    }],
                    "xAxes": [{
                        "ticks": {
                            "fontSize": 16
                        }
                    }]
                },
                "legend": {
                    "display": False
                }
            }
        }

        roi_path = self.output_dir / "chart_roi.png"
        if APIUtils.generate_quickchart(roi_config, str(roi_path), width=1200, height=600):
            charts['roi'] = roi_path
            print(f"  [OK] ROI chart saved")

        # Compounding Effect Chart
        compound_config = {
            "type": "line",
            "data": {
                "labels": ["Month 1", "Month 3", "Month 6", "Month 9", "Month 12"],
                "datasets": [{
                    "label": "Encoded Patterns",
                    "data": [50, 100, 150, 190, 221],
                    "borderColor": "rgb(0, 120, 212)",
                    "backgroundColor": "rgba(0, 120, 212, 0.1)",
                    "yAxisID": "y-axis-1",
                    "fill": True
                }, {
                    "label": "Monthly Value ($K)",
                    "data": [15, 22, 28, 35, 40],
                    "borderColor": "rgb(135, 100, 184)",
                    "backgroundColor": "rgba(135, 100, 184, 0.1)",
                    "yAxisID": "y-axis-2",
                    "fill": True
                }]
            },
            "options": {
                "title": {
                    "display": True,
                    "text": "The Compounding Effect",
                    "fontSize": 20
                },
                "scales": {
                    "yAxes": [{
                        "id": "y-axis-1",
                        "type": "linear",
                        "position": "left",
                        "ticks": {
                            "fontSize": 14
                        },
                        "scaleLabel": {
                            "display": True,
                            "labelString": "Encoded Patterns",
                            "fontSize": 14
                        }
                    }, {
                        "id": "y-axis-2",
                        "type": "linear",
                        "position": "right",
                        "ticks": {
                            "fontSize": 14
                        },
                        "scaleLabel": {
                            "display": True,
                            "labelString": "Monthly Value ($K)",
                            "fontSize": 14
                        },
                        "gridLines": {
                            "drawOnChartArea": False
                        }
                    }],
                    "xAxes": [{
                        "ticks": {
                            "fontSize": 14
                        }
                    }]
                }
            }
        }

        compound_path = self.output_dir / "chart_compounding.png"
        if APIUtils.generate_quickchart(compound_config, str(compound_path), width=1200, height=600):
            charts['compounding'] = compound_path
            print(f"  [OK] Compounding effect chart saved")

        return charts

    def generate_diagrams(self):
        """Generate diagrams using Mermaid"""
        print("\nGenerating diagrams...")

        diagrams = {}

        # Workers vs Tools Comparison
        comparison_mermaid = """
graph LR
    A[Tools ChatGPT] --> B[You start every task]
    A --> C[No memory stateless]
    A --> D[You provide context]
    A --> E[Suggestions only]

    F[Workers AI Workforce] --> G[Runs autonomously]
    F --> H[Remembers patterns]
    F --> I[Learns institutional knowledge]
    F --> J[Completed work product]

    style A fill:#ff6b6b
    style F fill:#4ecdc4
"""

        comparison_path = self.output_dir / "diagram_comparison.png"
        if APIUtils.generate_mermaid_diagram(comparison_mermaid, str(comparison_path)):
            diagrams['comparison'] = comparison_path
            print(f"  [OK] Comparison diagram saved")

        # Invoice Processor Workflow
        workflow_mermaid = """
graph TD
    A[Invoice arrives in email] --> B[Extract data]
    B --> C[Apply coding rules]
    C --> D{Amount > $50K?}
    D -->|Yes| E[Alert partner]
    D -->|No| F[Post to QuickBooks]
    E --> F
    F --> G[Log action]
    G --> H{Anomaly detected?}
    H -->|Yes| I[Flag for review]
    H -->|No| J[Complete]

    style A fill:#e3f2fd
    style J fill:#c8e6c9
    style I fill:#ffccbc
"""

        workflow_path = self.output_dir / "diagram_workflow.png"
        if APIUtils.generate_mermaid_diagram(workflow_mermaid, str(workflow_path)):
            diagrams['workflow'] = workflow_path
            print(f"  [OK] Workflow diagram saved")

        # Decision Tree
        decision_mermaid = """
graph TD
    A[Your Choice Today] --> B[Path A: Wait and See]
    A --> C[Path B: Start Encoding Now]
    B --> D[Competitors encode first]
    D --> E[Playing catch-up in 2027]
    C --> F[Build institutional capital]
    F --> G[Compounding advantage]

    style A fill:#fff3e0
    style B fill:#ffccbc
    style C fill:#c8e6c9
    style E fill:#ef5350
    style G fill:#66bb6a
"""

        decision_path = self.output_dir / "diagram_decision.png"
        if APIUtils.generate_mermaid_diagram(decision_mermaid, str(decision_path)):
            diagrams['decision'] = decision_path
            print(f"  [OK] Decision tree saved")

        return diagrams

    def build_presentation(self):
        """Build the complete presentation"""
        print("="*60)
        print("BUILDING AI WORKFORCE PRESENTATION")
        print("="*60)

        # Generate all assets
        historical_images = self.generate_historical_images()
        charts = self.generate_charts()
        diagrams = self.generate_diagrams()

        # Get paths to Figma illustrations
        illustrations = {
            'in_office': self.assets_dir / "illustrations" / "concepts" / "undraw_in_the_office_re_cxds.png",
            'creative_woman': self.assets_dir / "illustrations" / "concepts" / "undraw_creative_woman_re_u5tk.png",
            'in_thought': self.assets_dir / "illustrations" / "concepts" / "undraw_in_thought_re_qyxl.png",
            'team_collaboration': self.assets_dir / "illustrations" / "concepts" / "undraw_team_collaboration_re_ow29.png",
            'team_page': self.assets_dir / "illustrations" / "concepts" / "undraw_team_page_re_cffr.png",
            'our_solution': self.assets_dir / "illustrations" / "concepts" / "undraw_our_solution_re_8yk6.png",
            'startup_life': self.assets_dir / "illustrations" / "concepts" / "undraw_startup_life_re_8ow9.png",
            'building_blocks': self.assets_dir / "illustrations" / "concepts" / "undraw_building_blocks_re_5ahy.png",
            'interview': self.assets_dir / "illustrations" / "concepts" / "undraw_interview_re_e5jn.png",
            'swipe_options': self.assets_dir / "illustrations" / "concepts" / "undraw_swipe_profiles_re_tvqm.png",
            'live_collaboration': self.assets_dir / "illustrations" / "concepts" / "undraw_live_collaboration_re_60ha.png",
            'team_spirit': self.assets_dir / "illustrations" / "concepts" / "undraw_team_spirit_re_yl1v.png",
            'healthy_lifestyle': self.assets_dir / "illustrations" / "concepts" / "undraw_healthy_lifestyle_re_ifwg.png",
        }

        print("\n" + "="*60)
        print("CREATING SLIDES")
        print("="*60)

        # SLIDE 1: Title
        print("\n[1/28] Title slide")
        self.add_title_slide(
            "The AI Workforce",
            "Deploying Institutional Intelligence in Professional Services"
        )

        # ACT 1: HISTORICAL CONTEXT
        print("\n--- ACT 1: HISTORICAL CONTEXT ---")

        # SLIDE 2: The Harvest
        print("[2/28] The Harvest (1800s)")
        self.add_image_slide(
            "The Harvest (1800s)",
            historical_images.get('harvest', ''),
            "95% of Americans worked in agriculture"
        )

        # SLIDE 3: Combine Harvester
        print("[3/28] The Combine Harvester")
        self.add_image_slide(
            "The Combine Harvester Arrives",
            historical_images.get('combine', ''),
            "One machine eliminated 90% of harvest labor"
        )

        # SLIDE 4: Switchboard Operators
        print("[4/28] Switchboard Operators")
        self.add_image_slide(
            "Switchboard Operators (1950s)",
            historical_images.get('switchboard', ''),
            "Bell System employed 350,000 switchboard operators"
        )

        # SLIDE 5: Direct Dial
        print("[5/28] Direct Dial")
        self.add_image_slide(
            "Direct Dial Changes Everything",
            historical_images.get('direct_dial', ''),
            "Automatic switching eliminated 350,000 jobs - Workers moved to new telecommunications roles"
        )

        # ACT 2: THE CURRENT TRAP
        print("\n--- ACT 2: THE CURRENT TRAP ---")

        # SLIDE 6: Today's Reality
        print("[6/28] Today's Reality")
        self.add_content_slide(
            "Today's Professional Services Reality",
            [
                "Your most experienced people spend 60% of time on repetitive tasks",
                "",
                "• Senior accountant: invoice coding, data entry, report formatting",
                "• Lead attorney: document review, research memos, status updates",
                "• Top recruiter: resume screening, scheduling, follow-ups"
            ],
            illustrations.get('in_office')
        )

        # SLIDE 7: Million-Dollar Bottleneck
        print("[7/28] Million-Dollar Bottleneck")
        self.add_content_slide(
            "The Million-Dollar Bottleneck",
            [
                "$2.5M/year in rework costs",
                "",
                "18-24 months to rebuild patterns when key person leaves",
                "",
                "Tribal knowledge locked in email, Slack, individual brains"
            ],
            None
        )

        # SLIDE 8: The Tool Trap
        print("[8/28] The Tool Trap")
        self.add_content_slide(
            "The Tool Trap",
            [
                '"We bought ChatGPT licenses for everyone!"',
                "",
                "Result:",
                "• 12% adoption",
                "• Zero process change",
                "• Why? Tools require humans to operate them every time"
            ],
            illustrations.get('creative_woman')
        )

        # SLIDE 9: What If We Reframe?
        print("[9/28] What If We Reframe?")
        slide9 = self.add_content_slide(
            "What If We Reframe?",
            [
                "What if AI isn't a tool?",
                "",
                "What if it's a worker?"
            ],
            illustrations.get('in_thought')
        )
        self.add_morph_transition(slide9)

        # SLIDE 10: Workers vs Tools
        print("[10/28] Workers vs Tools")
        self.add_content_slide(
            "Workers vs. Tools",
            [
                "TOOLS (ChatGPT):",
                "• You start every task",
                "• No memory (stateless)",
                "• You provide context every time",
                "• Output: Suggestions",
                "",
                "WORKERS (AI Workforce):",
                "• Runs autonomously on triggers",
                "• Remembers patterns, history",
                "• Learns institutional knowledge",
                "• Output: Completed work product"
            ],
            None
        )

        # ACT 3: PROOF - MERIDIAN GROUP
        print("\n--- ACT 3: PROOF - MERIDIAN GROUP ---")

        # SLIDE 11: Meet Meridian
        print("[11/28] Meet Meridian Group")
        self.add_content_slide(
            "Meet Meridian Group",
            [
                "50-person accounting firm",
                "$10.5M annual revenue",
                "",
                "Typical professional services challenges:",
                "• High-value talent doing low-value work",
                "• Knowledge loss when people leave",
                "• Can't scale without adding headcount"
            ],
            illustrations.get('team_collaboration')
        )

        # SLIDE 12: 5 AI Workers
        print("[12/28] 5 AI Workers Deployed")
        self.add_content_slide(
            "They Deployed 5 AI Workers",
            [
                "1. Invoice Processor - Codes 800 invoices/month",
                "2. Pipeline Builder - Structures CRM data",
                "3. Outreach Drafter - Writes personalized client communications",
                "4. Talent Hunter - Screens candidates, schedules interviews",
                "5. Content Co-Pilot - Writes blog posts, social updates"
            ],
            illustrations.get('team_page')
        )

        # SLIDE 13: The Numbers
        print("[13/28] The Numbers")
        self.add_image_slide(
            "The Numbers",
            charts.get('roi', ''),
            None
        )

        # SLIDE 14: Invoice Processor Deep Dive
        print("[14/28] Invoice Processor Deep Dive")
        self.add_image_slide(
            "Invoice Processor - How It Works",
            diagrams.get('workflow', ''),
            "800 invoices/month • 120 hours saved • 99.4% accuracy"
        )

        # SLIDE 15: What Really Changed
        print("[15/28] What Really Changed")
        self.add_content_slide(
            "But Here's What Really Changed",
            [
                "221 institutional knowledge patterns encoded:",
                "",
                "• 'When client mentions year-end, create tax planning task'",
                "• 'When invoice >$50K, alert partner before posting'",
                "• 'When candidate mentions CPA, ask about state license'",
                "",
                "This is institutional capital being preserved"
            ],
            illustrations.get('our_solution')
        )

        # SLIDE 16: Three Autonomy Modes
        print("[16/28] Three Autonomy Modes")
        self.add_content_slide(
            "Three Autonomy Modes",
            [
                "1. AUTONOMOUS",
                "   Completes task, logs action",
                "   Example: Invoice coding",
                "",
                "2. SUPERVISED",
                "   Completes task, asks human to approve",
                "   Example: Client outreach",
                "",
                "3. HUMAN-LED",
                "   Drafts work, human edits/sends",
                "   Example: Blog posts"
            ],
            None
        )

        # SLIDE 17: What Happened to Humans?
        print("[17/28] What Happened to Humans?")
        self.add_content_slide(
            "What Happened to the Humans?",
            [
                "Didn't eliminate roles. Shifted focus to high-value work:",
                "",
                "• Senior accountant → client advisory",
                "  (billable hours +30%)",
                "",
                "• Recruiter → candidate relationships",
                "  (placements +40%)",
                "",
                "• Marketing manager → strategy",
                "  (campaigns +25%)"
            ],
            illustrations.get('startup_life')
        )

        # SLIDE 18: Compounding Effect
        print("[18/28] The Compounding Effect")
        self.add_image_slide(
            "The Compounding Effect",
            charts.get('compounding', ''),
            "Each pattern teaches the AI workers institutional judgment"
        )

        # ACT 4: THE CAPITAL LENS
        print("\n--- ACT 4: THE CAPITAL LENS ---")

        # SLIDE 19: The Real Asset
        print("[19/28] The Real Asset")
        self.add_content_slide(
            "The Real Asset",
            [
                "Your institutional knowledge is capital",
                "",
                "It's on balance sheet as 'goodwill'",
                "but locked in people's heads",
                "",
                "When they leave, you write it off"
            ],
            illustrations.get('building_blocks')
        )

        # SLIDE 20: Capital Lens Framework
        print("[20/28] Capital Lens Framework")
        self.add_content_slide(
            "The Capital Lens Framework",
            [
                "Financial ROI ←→ Institutional Knowledge Preservation",
                "",
                "$332,600             221 encoded patterns",
                "(Year 1)              (Permanent asset)",
                "",
                "",
                "Both matter, but knowledge compounds forever"
            ],
            None
        )

        # SLIDE 21: Scar Tissue Test
        print("[21/28] Scar Tissue Test")
        self.add_content_slide(
            "The Scar Tissue Test",
            [
                "How do you know where institutional knowledge lives?",
                "",
                "Find your scar tissue",
                "",
                "Ask: 'What mistake did we make that we'll never make again?'",
                "",
                "That's encoded institutional judgment"
            ],
            illustrations.get('interview')
        )

        # SLIDE 22: Examples from Meridian
        print("[22/28] Meridian Scar Tissue Examples")
        self.add_content_slide(
            "Examples from Meridian",
            [
                "1. 'Never send tax estimates on Fridays'",
                "   (learned after client panic calls)",
                "",
                "2. 'Always ask construction clients about prevailing wage'",
                "   (missed audit issue once)",
                "",
                "3. 'Flag invoices with manual journal entries'",
                "   (found fraud this way)",
                "",
                "Each one worth $10K-$100K in prevented losses"
            ],
            None
        )

        # ACT 5: THE DURABILITY QUESTION
        print("\n--- ACT 5: THE DURABILITY QUESTION ---")

        # SLIDE 23: But AI Is Changing
        print("[23/28] 'But AI Is Changing So Fast...'")
        slide23 = self.add_content_slide(
            '"But AI Is Changing So Fast..."',
            [
                "Common objection:",
                "",
                '"Should we wait for GPT-6?"',
                "",
                "This thinking mistakes the asset"
            ],
            illustrations.get('swipe_options')
        )
        self.add_morph_transition(slide23)

        # SLIDE 24: The Asset Isn't the AI Model
        print("[24/28] The Asset Isn't the AI Model")
        self.add_content_slide(
            "The Asset Isn't the AI Model",
            [
                "The asset: Your 221 encoded patterns",
                "",
                "The execution layer: Which AI model runs them",
                "",
                "When GPT-6 launches → upgrade execution layer",
                "",
                "Your patterns persist"
            ],
            illustrations.get('live_collaboration')
        )

        # SLIDE 25: Like Hiring
        print("[25/28] Like Hiring")
        self.add_content_slide(
            "Like Hiring",
            [
                "You hire people knowing they'll leave",
                "",
                "Value = work done while here + documentation they leave",
                "",
                "AI workers = same logic, but patterns persist forever",
                "",
                "No retirement, no job changes, no knowledge loss"
            ],
            illustrations.get('team_spirit')
        )

        # ACT 6: THE CALL
        print("\n--- ACT 6: THE CALL ---")

        # SLIDE 26: Two Paths Forward
        print("[26/28] Two Paths Forward")
        self.add_image_slide(
            "Two Paths Forward",
            diagrams.get('decision', ''),
            None
        )

        # SLIDE 27: Getting Started
        print("[27/28] Getting Started")
        self.add_content_slide(
            "Getting Started",
            [
                "1. Identify one scar tissue pattern (1 week)",
                "",
                "2. Deploy one AI worker (2-4 weeks)",
                "",
                "3. Measure, refine, encode next pattern (ongoing)",
                "",
                "4. Repeat until institutional knowledge is preserved (12-18 months)"
            ],
            illustrations.get('healthy_lifestyle')
        )

        # SLIDE 28: The Question
        print("[28/28] The Question")
        self.add_title_slide(
            "The harvest was eliminated.\nThe workers weren't.",
            ""
        )

        # Add final text to slide 28
        last_slide = self.prs.slides[-1]
        question_box = last_slide.shapes.add_textbox(
            Inches(1), Inches(6), Inches(14), Inches(2)
        )
        question_frame = question_box.text_frame
        question_frame.text = "What work will you eliminate?\nWhat will your workers do instead?"
        for para in question_frame.paragraphs:
            para.font.size = Pt(36)
            para.font.color.rgb = self.brand_blue
            para.alignment = PP_ALIGN.CENTER

        # Save presentation
        output_file = self.output_dir / "AI_Workforce_Live_Event.pptx"
        self.prs.save(str(output_file))

        print("\n" + "="*60)
        print("PRESENTATION COMPLETE")
        print("="*60)
        print(f"\nSaved to: {output_file}")
        print(f"Total slides: {len(self.prs.slides)}")
        print(f"\nAssets generated:")
        print(f"  Historical images: {len(historical_images)}")
        print(f"  Charts: {len(charts)}")
        print(f"  Diagrams: {len(diagrams)}")
        print("\nReady for rehearsal!")

        return output_file


def main():
    """Main execution"""
    generator = AIWorkforcePresentation()
    output_file = generator.build_presentation()
    print(f"\n[OK] Presentation ready: {output_file}")


if __name__ == "__main__":
    main()
